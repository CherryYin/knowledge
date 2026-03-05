"""
文档解析模块 (Document Parser)

支持格式: PDF (via MinerU), DOCX, PPTX, XLSX/CSV, MD, HTML, TXT, 图片, LaTeX, EPUB
输出: 统一的 ParsedDocument 结构，包含文本、表格、公式、图片块
"""

import hashlib
import os
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Optional


# ─────────────────────────────────────────────
#  工具函数
# ─────────────────────────────────────────────

def generate_doc_id(file_path: str) -> str:
    """根据文件路径 + 文件内容生成唯一文档 ID"""
    h = hashlib.md5(file_path.encode())
    try:
        stat = os.stat(file_path)
        h.update(str(stat.st_size).encode())
        h.update(str(stat.st_mtime).encode())
    except OSError:
        pass
    return h.hexdigest()


# ─────────────────────────────────────────────
#  数据模型
# ─────────────────────────────────────────────

class ContentType(Enum):
    TEXT = "text"
    TABLE = "table"
    FORMULA = "formula"
    IMAGE = "image"
    CODE = "code"
    HEADING = "heading"


@dataclass
class ContentBlock:
    """统一内容块"""
    type: ContentType
    content: str                                      # 文本内容或文件路径
    page: Optional[int] = None                        # 所在页码
    bbox: Optional[list] = None                       # 位置坐标 [x0, y0, x1, y1]
    metadata: dict = field(default_factory=dict)      # 额外元信息
    image_path: Optional[str] = None                  # 关联的图片路径（表格截图等）
    children: list = field(default_factory=list)      # 子元素


@dataclass
class ParsedDocument:
    """解析后的文档"""
    doc_id: str
    filename: str
    format: str
    title: Optional[str] = None
    author: Optional[str] = None
    page_count: Optional[int] = None
    blocks: list = field(default_factory=list)        # list[ContentBlock]
    raw_text: str = ""
    metadata: dict = field(default_factory=dict)


# ─────────────────────────────────────────────
#  MinerU PDF 解析引擎
# ─────────────────────────────────────────────

class MinerUParser:
    """
    基于 MinerU (magic-pdf) 的 PDF 深度解析器
    内置 Layout 分析、表格识别、公式识别、OCR

    安装: pip install magic-pdf[full]
    """

    def __init__(self, use_gpu: bool = True):
        self.use_gpu = use_gpu

    def parse_pdf(self, pdf_path: str, output_dir: str) -> dict:
        """
        解析 PDF，输出结构化内容

        返回:
            {
                "text_blocks": [...],
                "tables": [...],
                "formulas": [...],
                "images": [...],
            }
        """
        try:
            from magic_pdf.data.data_reader_writer import (
                FileBasedDataWriter,
                FileBasedDataReader,
            )
            from magic_pdf.pipe.UNIPipe import UNIPipe
        except ImportError as exc:
            raise ImportError(
                "magic-pdf is not installed. Run: pip install magic-pdf[full]"
            ) from exc

        reader = FileBasedDataReader("")
        pdf_bytes = reader.read(pdf_path)
        writer = FileBasedDataWriter(output_dir)

        pipe = UNIPipe(pdf_bytes, jso_useful_key={}, image_writer=writer)
        pipe.pipe_classify()
        pipe.pipe_analyze()
        pipe.pipe_parse()

        content_list = pipe.pipe_mk_uni_format(output_dir)
        pipe.pipe_mk_markdown(output_dir)

        return self._structure_output(content_list, output_dir)

    def _structure_output(self, content_list: list, output_dir: str) -> dict:
        """将 MinerU 输出转化为统一结构"""
        result: dict = {
            "text_blocks": [],
            "tables": [],
            "formulas": [],
            "images": [],
        }

        for block in content_list:
            block_type = block.get("type")

            if block_type == "text":
                result["text_blocks"].append({
                    "content": block.get("text", ""),
                    "page": block.get("page_idx"),
                    "bbox": block.get("bbox"),
                    "is_title": block.get("is_title", False),
                    "level": block.get("level", 0),
                })

            elif block_type == "table":
                result["tables"].append({
                    "html": block.get("html", ""),
                    "markdown": block.get("markdown", ""),
                    "cells": block.get("cells", []),
                    "page": block.get("page_idx"),
                    "bbox": block.get("bbox"),
                    "image_path": block.get("img_path"),
                })

            elif block_type == "equation":
                result["formulas"].append({
                    "latex": block.get("latex", ""),
                    "page": block.get("page_idx"),
                    "bbox": block.get("bbox"),
                    "inline": block.get("inline", False),
                    "image_path": block.get("img_path"),
                })

            elif block_type == "image":
                result["images"].append({
                    "path": block.get("img_path"),
                    "caption": block.get("caption", ""),
                    "page": block.get("page_idx"),
                    "bbox": block.get("bbox"),
                })

        return result


# ─────────────────────────────────────────────
#  解析器基类与各格式实现
# ─────────────────────────────────────────────

class BaseParser(ABC):
    """解析器基类"""

    @abstractmethod
    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        pass

    @abstractmethod
    def supported_formats(self) -> list:
        pass


class PDFParser(BaseParser):
    """PDF 解析器 — 基于 MinerU"""

    def __init__(self, use_gpu: bool = True):
        self.mineru = MinerUParser(use_gpu=use_gpu)

    def supported_formats(self) -> list:
        return [".pdf"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        result = self.mineru.parse_pdf(file_path, output_dir)
        blocks = []

        for tb in result["text_blocks"]:
            btype = ContentType.HEADING if tb["is_title"] else ContentType.TEXT
            blocks.append(ContentBlock(
                type=btype,
                content=tb["content"],
                page=tb["page"],
                bbox=tb["bbox"],
                metadata={"level": tb.get("level", 0)},
            ))

        for table in result["tables"]:
            blocks.append(ContentBlock(
                type=ContentType.TABLE,
                content=table["markdown"],
                page=table["page"],
                bbox=table["bbox"],
                image_path=table.get("image_path"),
                metadata={"html": table["html"]},
            ))

        for formula in result["formulas"]:
            blocks.append(ContentBlock(
                type=ContentType.FORMULA,
                content=formula["latex"],
                page=formula["page"],
                bbox=formula["bbox"],
                image_path=formula.get("image_path"),
                metadata={"inline": formula["inline"]},
            ))

        for img in result["images"]:
            blocks.append(ContentBlock(
                type=ContentType.IMAGE,
                content=img["caption"],
                page=img["page"],
                bbox=img["bbox"],
                image_path=img["path"],
            ))

        # 按页码 + y 坐标排序
        blocks.sort(key=lambda b: (b.page or 0, (b.bbox or [0, 0])[1]))

        raw_text = "\n".join(
            b.content for b in blocks if b.type in (ContentType.TEXT, ContentType.HEADING)
        )

        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="pdf",
            blocks=blocks,
            raw_text=raw_text,
        )


class DocxParser(BaseParser):
    """Word 文档解析器 — 基于 python-docx"""

    def supported_formats(self) -> list:
        return [".docx", ".doc"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        try:
            from docx import Document as DocxDocument
        except ImportError as exc:
            raise ImportError("python-docx is not installed. Run: pip install python-docx") from exc

        doc = DocxDocument(file_path)
        blocks = []

        for element in doc.element.body:
            tag = element.tag.split("}")[-1]

            if tag == "p":
                para = self._find_paragraph(doc, element)
                if para and para.text.strip():
                    style = para.style.name if para.style else ""
                    is_heading = style.lower().startswith("heading")
                    blocks.append(ContentBlock(
                        type=ContentType.HEADING if is_heading else ContentType.TEXT,
                        content=para.text.strip(),
                        metadata={"style": style},
                    ))

            elif tag == "tbl":
                table = self._find_table(doc, element)
                if table:
                    md = self._table_to_markdown(table)
                    blocks.append(ContentBlock(
                        type=ContentType.TABLE,
                        content=md,
                        metadata={"rows": len(table.rows), "cols": len(table.columns)},
                    ))

        # 提取图片
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                img_data = rel.target_part.blob
                img_filename = Path(rel.target_ref).name
                img_path = Path(output_dir) / img_filename
                img_path.parent.mkdir(parents=True, exist_ok=True)
                img_path.write_bytes(img_data)
                blocks.append(ContentBlock(
                    type=ContentType.IMAGE,
                    content="",
                    image_path=str(img_path),
                ))

        raw_text = "\n".join(
            b.content for b in blocks if b.type in (ContentType.TEXT, ContentType.HEADING)
        )
        props = doc.core_properties

        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="docx",
            title=props.title or "",
            author=props.author or "",
            blocks=blocks,
            raw_text=raw_text,
        )

    def _find_paragraph(self, doc, element):
        for para in doc.paragraphs:
            if para._element is element:
                return para
        return None

    def _find_table(self, doc, element):
        for table in doc.tables:
            if table._element is element:
                return table
        return None

    def _table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if len(rows) > 1:
            sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
            rows.insert(1, sep)
        return "\n".join(rows)


class PptxParser(BaseParser):
    """PowerPoint 解析器 — 基于 python-pptx"""

    def supported_formats(self) -> list:
        return [".pptx", ".ppt"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError as exc:
            raise ImportError("python-pptx is not installed. Run: pip install python-pptx") from exc

        prs = Presentation(file_path)
        blocks = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        # 首个占位符（标题框）视为标题
                        is_title = shape.shape_type == 13 or (
                            hasattr(shape, "placeholder_format")
                            and shape.placeholder_format is not None
                            and shape.placeholder_format.idx == 0
                        )
                        blocks.append(ContentBlock(
                            type=ContentType.HEADING if is_title else ContentType.TEXT,
                            content=text,
                            page=slide_idx,
                            metadata={"slide": slide_idx},
                        ))

                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        img_bytes = shape.image.blob
                        ext = shape.image.ext
                        img_name = f"slide{slide_idx}_img{shape.shape_id}.{ext}"
                        img_path = Path(output_dir) / img_name
                        img_path.parent.mkdir(parents=True, exist_ok=True)
                        img_path.write_bytes(img_bytes)
                        blocks.append(ContentBlock(
                            type=ContentType.IMAGE,
                            content="",
                            page=slide_idx,
                            image_path=str(img_path),
                        ))
                    except Exception:
                        pass

        raw_text = "\n".join(
            b.content for b in blocks if b.type in (ContentType.TEXT, ContentType.HEADING)
        )
        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="pptx",
            blocks=blocks,
            raw_text=raw_text,
        )


class SpreadsheetParser(BaseParser):
    """Excel / CSV 解析器 — 基于 pandas"""

    def supported_formats(self) -> list:
        return [".xlsx", ".xls", ".csv"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        try:
            import pandas as pd
        except ImportError as exc:
            raise ImportError("pandas is not installed. Run: pip install pandas openpyxl") from exc

        ext = Path(file_path).suffix.lower()
        blocks = []

        if ext == ".csv":
            df_dict = {"Sheet1": pd.read_csv(file_path)}
        else:
            df_dict = pd.read_excel(file_path, sheet_name=None)

        for sheet_name, df in df_dict.items():
            df = df.fillna("")
            md = df.to_markdown(index=False)
            blocks.append(ContentBlock(
                type=ContentType.TABLE,
                content=md,
                metadata={"sheet": sheet_name, "rows": len(df), "cols": len(df.columns)},
            ))

        raw_text = "\n\n".join(b.content for b in blocks)
        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format=ext.lstrip("."),
            blocks=blocks,
            raw_text=raw_text,
        )


class MarkdownParser(BaseParser):
    """Markdown 解析器"""

    def supported_formats(self) -> list:
        return [".md", ".markdown"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        import re

        text = Path(file_path).read_text(encoding="utf-8")
        blocks = []

        for line in text.splitlines():
            heading_match = re.match(r"^(#{1,6})\s+(.*)", line)
            if heading_match:
                level = len(heading_match.group(1))
                blocks.append(ContentBlock(
                    type=ContentType.HEADING,
                    content=heading_match.group(2).strip(),
                    metadata={"level": level},
                ))
            elif line.strip():
                blocks.append(ContentBlock(
                    type=ContentType.TEXT,
                    content=line.strip(),
                ))

        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="md",
            blocks=blocks,
            raw_text=text,
        )


class HTMLParser(BaseParser):
    """HTML 解析器 — 基于 BeautifulSoup4"""

    def supported_formats(self) -> list:
        return [".html", ".htm"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            raise ImportError("beautifulsoup4 is not installed. Run: pip install beautifulsoup4 lxml") from exc

        html = Path(file_path).read_text(encoding="utf-8")
        soup = BeautifulSoup(html, "lxml")
        blocks = []

        for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "table"]):
            name = tag.name
            if name in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(name[1])
                text = tag.get_text(strip=True)
                if text:
                    blocks.append(ContentBlock(
                        type=ContentType.HEADING,
                        content=text,
                        metadata={"level": level},
                    ))
            elif name == "p":
                text = tag.get_text(strip=True)
                if text:
                    blocks.append(ContentBlock(type=ContentType.TEXT, content=text))
            elif name == "table":
                rows = []
                for tr in tag.find_all("tr"):
                    cells = [td.get_text(strip=True).replace("|", "\\|") for td in tr.find_all(["td", "th"])]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    sep = "| " + " | ".join(["---"] * len(rows[0].split("|"))) + " |"
                    rows.insert(1, sep)
                    blocks.append(ContentBlock(type=ContentType.TABLE, content="\n".join(rows)))

        raw_text = "\n".join(
            b.content for b in blocks if b.type in (ContentType.TEXT, ContentType.HEADING)
        )
        title_tag = soup.find("title")
        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="html",
            title=title_tag.get_text(strip=True) if title_tag else "",
            blocks=blocks,
            raw_text=raw_text,
        )


class TextParser(BaseParser):
    """纯文本解析器"""

    def supported_formats(self) -> list:
        return [".txt", ".text"]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        text = Path(file_path).read_text(encoding="utf-8", errors="replace")
        blocks = [ContentBlock(type=ContentType.TEXT, content=para.strip())
                  for para in text.split("\n\n") if para.strip()]
        return ParsedDocument(
            doc_id=generate_doc_id(file_path),
            filename=Path(file_path).name,
            format="txt",
            blocks=blocks,
            raw_text=text,
        )


# ─────────────────────────────────────────────
#  格式路由器
# ─────────────────────────────────────────────

class FormatRouter:
    """
    格式路由器 — 根据文件扩展名选择合适的解析器

    默认已注册: PDF, DOCX, PPTX, XLSX/CSV, MD, HTML, TXT
    """

    def __init__(self, use_gpu: bool = True):
        self._parsers: dict = {}
        # 注册默认解析器
        self.register(PDFParser(use_gpu=use_gpu))
        self.register(DocxParser())
        self.register(PptxParser())
        self.register(SpreadsheetParser())
        self.register(MarkdownParser())
        self.register(HTMLParser())
        self.register(TextParser())

    def register(self, parser: BaseParser) -> None:
        for fmt in parser.supported_formats():
            self._parsers[fmt.lower()] = parser

    def route(self, file_path: str) -> BaseParser:
        ext = Path(file_path).suffix.lower()
        if ext == ".doc":
            ext = ".docx"
        if ext not in self._parsers:
            raise ValueError(f"Unsupported file format: {ext!r}")
        return self._parsers[ext]

    def parse(self, file_path: str, output_dir: str) -> ParsedDocument:
        """解析文件，返回统一的 ParsedDocument"""
        parser = self.route(file_path)
        return parser.parse(file_path, output_dir)

    @property
    def supported_extensions(self) -> list:
        return sorted(self._parsers.keys())
